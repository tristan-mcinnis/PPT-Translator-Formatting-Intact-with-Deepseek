"""PowerPoint translation pipeline utilities."""
from __future__ import annotations

import json
import xml.etree.ElementTree as ET
from concurrent.futures import ThreadPoolExecutor
from pathlib import Path
from xml.dom import minidom

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

from .translation import TranslationService


def get_alignment_value(alignment_str: str | None):
    """Convert alignment string to PP_ALIGN enum value."""
    alignment_map = {
        "PP_ALIGN.CENTER": PP_ALIGN.CENTER,
        "PP_ALIGN.LEFT": PP_ALIGN.LEFT,
        "PP_ALIGN.RIGHT": PP_ALIGN.RIGHT,
        "PP_ALIGN.JUSTIFY": PP_ALIGN.JUSTIFY,
        "None": None,
        None: None,
    }
    return alignment_map.get(alignment_str)


def get_shape_properties(shape):
    """Extract text shape properties."""
    shape_data = {
        "text": "",
        "font_size": None,
        "font_name": None,
        "alignment": None,
        "width": shape.width,
        "height": shape.height,
        "left": shape.left,
        "top": shape.top,
        "bold": None,
        "italic": None,
        "line_spacing": None,
        "space_before": None,
        "space_after": None,
        "font_color": None,
    }
    if hasattr(shape, "text"):
        shape_data["text"] = shape.text.strip()
        if hasattr(shape, "text_frame"):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if getattr(run.font, "size", None) is not None:
                        shape_data["font_size"] = run.font.size.pt
                    if getattr(run.font, "name", None):
                        shape_data["font_name"] = run.font.name
                    if hasattr(run.font, "bold"):
                        shape_data["bold"] = run.font.bold
                    if hasattr(run.font, "italic"):
                        shape_data["italic"] = run.font.italic
                    if (
                        getattr(run.font, "color", None) is not None
                        and getattr(run.font.color, "rgb", None) is not None
                    ):
                        shape_data["font_color"] = str(run.font.color.rgb)
                if getattr(paragraph, "line_spacing", None) is not None:
                    shape_data["line_spacing"] = paragraph.line_spacing
                if getattr(paragraph, "space_before", None) is not None:
                    shape_data["space_before"] = paragraph.space_before
                if getattr(paragraph, "space_after", None) is not None:
                    shape_data["space_after"] = paragraph.space_after
                if getattr(paragraph, "alignment", None) is not None:
                    shape_data["alignment"] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None
    return shape_data


def apply_shape_properties(shape, shape_data):
    """Apply saved properties to a shape."""
    try:
        shape.width = shape_data["width"]
        shape.height = shape_data["height"]
        shape.left = shape_data["left"]
        shape.top = shape_data["top"]
        shape.text = ""
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = shape_data["text"]
        if shape_data.get("font_size"):
            adjusted_size = shape_data["font_size"] * 0.7
            run.font.size = Pt(adjusted_size)
        run.font.name = shape_data.get("font_name") or "Arial"
        if shape_data.get("font_color"):
            run.font.color.rgb = RGBColor.from_string(shape_data["font_color"])
        if shape_data.get("bold") is not None:
            run.font.bold = shape_data["bold"]
        if shape_data.get("italic") is not None:
            run.font.italic = shape_data["italic"]
        if shape_data.get("alignment"):
            paragraph.alignment = get_alignment_value(shape_data["alignment"])
        if shape_data.get("line_spacing"):
            paragraph.line_spacing = shape_data["line_spacing"]
        if shape_data.get("space_before"):
            paragraph.space_before = shape_data["space_before"]
        if shape_data.get("space_after"):
            paragraph.space_after = shape_data["space_after"]
    except Exception as exc:  # pragma: no cover - best effort logging
        print(f"Error applying shape properties: {exc}")


def get_table_properties(table):
    """Extract table properties."""
    table_data = {
        "rows": len(table.rows),
        "cols": len(table.columns),
        "cells": [],
    }
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            cell_data = {
                "text": cell.text.strip(),
                "font_size": None,
                "font_name": None,
                "alignment": None,
                "margin_left": cell.margin_left,
                "margin_right": cell.margin_right,
                "margin_top": cell.margin_top,
                "margin_bottom": cell.margin_bottom,
                "vertical_anchor": str(cell.vertical_anchor) if cell.vertical_anchor else None,
                "font_color": None,
            }
            if cell.text_frame.paragraphs:
                paragraph = cell.text_frame.paragraphs[0]
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if getattr(run.font, "size", None) is not None:
                        cell_data["font_size"] = run.font.size.pt
                    if getattr(run.font, "name", None):
                        cell_data["font_name"] = run.font.name
                    if hasattr(run.font, "bold"):
                        cell_data["bold"] = run.font.bold
                    if hasattr(run.font, "italic"):
                        cell_data["italic"] = run.font.italic
                    if (
                        getattr(run.font, "color", None) is not None
                        and getattr(run.font.color, "rgb", None) is not None
                    ):
                        cell_data["font_color"] = str(run.font.color.rgb)
                if getattr(paragraph, "alignment", None) is not None:
                    cell_data["alignment"] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None
            row_data.append(cell_data)
        table_data["cells"].append(row_data)
    return table_data


def apply_table_properties(table, table_data):
    """Apply saved table properties."""
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            try:
                cell_data = table_data["cells"][row_idx][col_idx]
                cell.margin_left = cell_data["margin_left"]
                cell.margin_right = cell_data["margin_right"]
                cell.margin_top = cell_data["margin_top"]
                cell.margin_bottom = cell_data["margin_bottom"]
                if cell_data.get("vertical_anchor"):
                    cell.vertical_anchor = eval(cell_data["vertical_anchor"])
                cell.text = ""
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = cell_data["text"]
                if cell_data.get("font_size"):
                    adjusted_size = cell_data["font_size"] * 0.8
                    run.font.size = Pt(adjusted_size)
                run.font.name = cell_data.get("font_name") or "Arial"
                if cell_data.get("font_color"):
                    run.font.color.rgb = RGBColor.from_string(cell_data["font_color"])
                if "bold" in cell_data:
                    run.font.bold = cell_data["bold"]
                if "italic" in cell_data:
                    run.font.italic = cell_data["italic"]
                if cell_data.get("alignment"):
                    paragraph.alignment = get_alignment_value(cell_data["alignment"])
            except Exception as exc:  # pragma: no cover - best effort logging
                print(f"Error setting cell properties: {exc}")


def extract_text_from_slide(
    slide,
    slide_number: int,
    *,
    translator: TranslationService | None,
    source_lang: str,
    target_lang: str,
):
    """Extract text from a slide and optionally translate it."""
    slide_element = ET.Element("slide")
    slide_element.set("number", str(slide_number))
    for shape_index, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_element = ET.SubElement(slide_element, "table_element")
            table_element.set("shape_index", str(shape_index))
            table_data = get_table_properties(shape.table)
            if translator:
                for row in table_data["cells"]:
                    for cell in row:
                        cell["text"] = translator.translate(cell["text"], source_lang, target_lang)
            props_element = ET.SubElement(table_element, "properties")
            props_element.text = json.dumps(table_data, indent=2)
        elif hasattr(shape, "text"):
            text_element = ET.SubElement(slide_element, "text_element")
            text_element.set("shape_index", str(shape_index))
            shape_data = get_shape_properties(shape)
            if translator:
                shape_data["text"] = translator.translate(shape_data["text"], source_lang, target_lang)
            props_element = ET.SubElement(text_element, "properties")
            props_element.text = json.dumps(shape_data, indent=2)
    return slide_element


def ppt_to_xml(
    ppt_path: str,
    *,
    translator: TranslationService | None,
    source_lang: str,
    target_lang: str,
    max_workers: int = 4,
) -> Optional[str]:
    """Convert a PowerPoint presentation to XML."""
    root = ET.Element("presentation")
    base_dir = Path(ppt_path).parent
    try:
        prs = Presentation(ppt_path)
        root.set("file_path", Path(ppt_path).name)
        workers = max(1, max_workers)
        with ThreadPoolExecutor(max_workers=workers) as executor:
            future_to_slide = {
                executor.submit(
                    extract_text_from_slide,
                    slide,
                    slide_number,
                    translator=translator,
                    source_lang=source_lang,
                    target_lang=target_lang,
                ): slide_number
                for slide_number, slide in enumerate(prs.slides, start=1)
            }
            for future, slide_number in future_to_slide.items():
                slide_element = future.result()
                root.append(slide_element)
                intermediate_path = base_dir / f"slide_{slide_number}_{'translated' if translator else 'original'}.xml"
                xml_str = minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
                with open(intermediate_path, "w", encoding="utf-8") as handle:
                    handle.write(xml_str)
        return minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
    except Exception as exc:  # pragma: no cover - best effort logging
        print(f"Error processing presentation: {exc}")
        return None


def create_translated_ppt(original_ppt_path: str, translated_xml_path: str, output_ppt_path: str) -> None:
    """Create a new PowerPoint presentation using translated content."""
    try:
        prs = Presentation(original_ppt_path)
        tree = ET.parse(translated_xml_path)
        root = tree.getroot()
        for slide_number, slide in enumerate(prs.slides, start=1):
            xml_slide = root.find(f".//slide[@number='{slide_number}']")
            if xml_slide is None:
                continue
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_element = xml_slide.find(f".//table_element[@shape_index='{shape_index}']")
                    if table_element is not None:
                        props_element = table_element.find("properties")
                        if props_element is not None and props_element.text:
                            try:
                                table_data = json.loads(props_element.text)
                                apply_table_properties(shape.table, table_data)
                            except Exception as exc:  # pragma: no cover
                                print(f"Error applying table properties: {exc}")
                elif hasattr(shape, "text"):
                    text_element = xml_slide.find(f".//text_element[@shape_index='{shape_index}']")
                    if text_element is not None:
                        props_element = text_element.find("properties")
                        if props_element is not None and props_element.text:
                            try:
                                shape_data = json.loads(props_element.text)
                                apply_shape_properties(shape, shape_data)
                            except Exception as exc:  # pragma: no cover
                                print(f"Error applying shape properties: {exc}")
        prs.save(output_ppt_path)
        print(f"Translated PowerPoint saved to: {output_ppt_path}")
    except Exception as exc:  # pragma: no cover - logging only
        print(f"Error creating translated PowerPoint: {exc}")


def cleanup_intermediate_files(base_dir: Path, pattern: str = "slide_*.xml") -> None:
    """Remove intermediate XML files."""
    try:
        for file in base_dir.glob(pattern):
            file.unlink()
    except Exception as exc:  # pragma: no cover - logging only
        print(f"Warning: Could not clean up intermediate files: {exc}")


def process_ppt_file(
    ppt_path: Path,
    *,
    translator: TranslationService,
    source_lang: str,
    target_lang: str,
    max_workers: int = 4,
    cleanup: bool = True,
) -> Optional[Path]:
    """Process a single PowerPoint file from extraction to translated output."""
    if not ppt_path.is_file():
        raise FileNotFoundError(f"'{ppt_path}' is not a valid file.")
    if ppt_path.suffix.lower() not in {".ppt", ".pptx"}:
        raise ValueError(f"'{ppt_path}' is not a PowerPoint file.")

    base_dir = ppt_path.parent

    print(f"Generating original XML for {ppt_path.name}...")
    original_xml = ppt_to_xml(
        str(ppt_path),
        translator=None,
        source_lang=source_lang,
        target_lang=target_lang,
        max_workers=max_workers,
    )
    if original_xml:
        original_output_path = base_dir / f"{ppt_path.stem}_original.xml"
        with open(original_output_path, "w", encoding="utf-8") as handle:
            handle.write(original_xml)
        print(f"Original XML saved: {original_output_path}")

    print(
        f"Generating translated XML (from {source_lang} to {target_lang}) for {ppt_path.name}..."
    )
    translated_xml = ppt_to_xml(
        str(ppt_path),
        translator=translator,
        source_lang=source_lang,
        target_lang=target_lang,
        max_workers=max_workers,
    )
    if not translated_xml:
        return None

    translated_output_path = base_dir / f"{ppt_path.stem}_translated.xml"
    with open(translated_output_path, "w", encoding="utf-8") as handle:
        handle.write(translated_xml)
    print(f"Translated XML saved: {translated_output_path}")

    print(f"Creating translated PPT for {ppt_path.name}...")
    output_filename = f"{ppt_path.stem}_translated{ppt_path.suffix}"
    output_ppt_path = base_dir / output_filename
    create_translated_ppt(str(ppt_path), str(translated_output_path), str(output_ppt_path))

    if cleanup:
        cleanup_intermediate_files(base_dir)
        print("Cleanup complete.")

    return output_ppt_path
