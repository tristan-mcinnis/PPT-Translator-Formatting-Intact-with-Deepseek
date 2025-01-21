from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor
import xml.etree.ElementTree as ET
from xml.dom import minidom
import argparse
import os
import json
from openai import OpenAI
from dotenv import load_dotenv
from pathlib import Path
import sys
from concurrent.futures import ThreadPoolExecutor
from typing import Dict, Optional
import threading

# Load environment variables
load_dotenv()

# Initialize DeepSeek client
client = OpenAI(
    api_key=os.getenv("DEEPSEEK_API_KEY"),
    base_url="https://api.deepseek.com"
)

# Thread-safe translation cache
translation_cache: Dict[str, str] = {}
cache_lock = threading.Lock()

def translate_text(text: str, source_lang: str = 'zh', target_lang: str = 'en') -> str:
    """Translate text from source_lang to target_lang using DeepSeek with chunking."""
    try:
        if not text or text.isspace():
            return text
        chunks = chunk_text(text)
        translated_chunks = []
        for chunk in chunks:
            response = client.chat.completions.create(
                model="deepseek-chat",
                messages=[
                    {"role": "system", "content": f"You are a translator. Translate the following {source_lang} text to {target_lang}. Keep the formatting and maintain a natural, fluent translation."},
                    {"role": "user", "content": chunk}
                ],
                temperature=0.3,
                stream=False
            )
            translated_chunks.append(response.choices[0].message.content)
        return ' '.join(translated_chunks)
    except Exception as e:
        print(f"Translation error: {str(e)}")
        return text

def get_alignment_value(alignment_str):
    """Convert alignment string to PP_ALIGN enum value."""
    alignment_map = {
        'PP_ALIGN.CENTER': PP_ALIGN.CENTER,
        'PP_ALIGN.LEFT': PP_ALIGN.LEFT,
        'PP_ALIGN.RIGHT': PP_ALIGN.RIGHT,
        'PP_ALIGN.JUSTIFY': PP_ALIGN.JUSTIFY,
        'None': None
    }
    return alignment_map.get(alignment_str)

def get_shape_properties(shape):
    """Extract all properties from a shape."""
    shape_data = {
        'text': '',
        'font_size': None,
        'font_name': None,
        'alignment': None,
        'width': shape.width,
        'height': shape.height,
        'left': shape.left,
        'top': shape.top,
        'bold': None,
        'italic': None,
        'line_spacing': None,
        'space_before': None,
        'space_after': None,
        'font_color': None
    }
    if hasattr(shape, "text"):
        shape_data['text'] = shape.text.strip()
        if hasattr(shape, 'text_frame'):
            for paragraph in shape.text_frame.paragraphs:
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if hasattr(run.font, 'size') and run.font.size is not None:
                        shape_data['font_size'] = run.font.size.pt
                    if hasattr(run.font, 'name'):
                        shape_data['font_name'] = run.font.name
                    if hasattr(run.font, 'bold'):
                        shape_data['bold'] = run.font.bold
                    if hasattr(run.font, 'italic'):
                        shape_data['italic'] = run.font.italic
                    if (hasattr(run.font, 'color') and 
                        run.font.color is not None and 
                        hasattr(run.font.color, 'rgb') and 
                        run.font.color.rgb is not None):
                        shape_data['font_color'] = str(run.font.color.rgb)
                if hasattr(paragraph, 'line_spacing'):
                    shape_data['line_spacing'] = paragraph.line_spacing
                if hasattr(paragraph, 'space_before'):
                    shape_data['space_before'] = paragraph.space_before
                if hasattr(paragraph, 'space_after'):
                    shape_data['space_after'] = paragraph.space_after
                if hasattr(paragraph, 'alignment'):
                    shape_data['alignment'] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None
    return shape_data

def apply_shape_properties(shape, shape_data):
    """Apply saved properties to a shape."""
    try:
        shape.width = shape_data['width']
        shape.height = shape_data['height']
        shape.left = shape_data['left']
        shape.top = shape_data['top']
        shape.text = ""
        paragraph = shape.text_frame.paragraphs[0]
        run = paragraph.add_run()
        run.text = shape_data['text']
        if shape_data['font_size']:
            # reduce original font size by 30%
            adjusted_size = shape_data['font_size'] * 0.7
            run.font.size = Pt(adjusted_size)
        run.font.name = 'Arial'
        if shape_data.get('font_color'):
            run.font.color.rgb = RGBColor.from_string(shape_data['font_color'])
        if shape_data['bold'] is not None:
            run.font.bold = shape_data['bold']
        if shape_data['italic'] is not None:
            run.font.italic = shape_data['italic']
        if shape_data['alignment']:
            paragraph.alignment = get_alignment_value(shape_data['alignment'])
        if shape_data['line_spacing']:
            paragraph.line_spacing = shape_data['line_spacing']
        if shape_data['space_before']:
            paragraph.space_before = shape_data['space_before']
        if shape_data['space_after']:
            paragraph.space_after = shape_data['space_after']
    except Exception as e:
        print(f"Error applying shape properties: {str(e)}")

def get_table_properties(table):
    """Extract complete table properties."""
    table_data = {
        'rows': len(table.rows),
        'cols': len(table.columns),
        'cells': []
    }
    for row in table.rows:
        row_data = []
        for cell in row.cells:
            cell_data = {
                'text': cell.text.strip(),
                'font_size': None,
                'font_name': None,
                'alignment': None,
                'margin_left': cell.margin_left,
                'margin_right': cell.margin_right,
                'margin_top': cell.margin_top,
                'margin_bottom': cell.margin_bottom,
                'vertical_anchor': str(cell.vertical_anchor) if cell.vertical_anchor else None,
                'font_color': None
            }
            if cell.text_frame.paragraphs:
                paragraph = cell.text_frame.paragraphs[0]
                if paragraph.runs:
                    run = paragraph.runs[0]
                    if hasattr(run.font, 'size') and run.font.size is not None:
                        cell_data['font_size'] = run.font.size.pt
                    if hasattr(run.font, 'name'):
                        cell_data['font_name'] = run.font.name
                    if hasattr(run.font, 'bold'):
                        cell_data['bold'] = run.font.bold
                    if hasattr(run.font, 'italic'):
                        cell_data['italic'] = run.font.italic
                    if (hasattr(run.font, 'color') and 
                        run.font.color is not None and 
                        hasattr(run.font.color, 'rgb') and 
                        run.font.color.rgb is not None):
                        cell_data['font_color'] = str(run.font.color.rgb)
                if hasattr(paragraph, 'alignment'):
                    cell_data['alignment'] = f"PP_ALIGN.{paragraph.alignment}" if paragraph.alignment else None
            row_data.append(cell_data)
        table_data['cells'].append(row_data)
    return table_data

def apply_table_properties(table, table_data):
    """Apply saved table properties to a table."""
    for row_idx, row in enumerate(table.rows):
        for col_idx, cell in enumerate(row.cells):
            try:
                cell_data = table_data['cells'][row_idx][col_idx]
                cell.margin_left = cell_data['margin_left']
                cell.margin_right = cell_data['margin_right']
                cell.margin_top = cell_data['margin_top']
                cell.margin_bottom = cell_data['margin_bottom']
                if cell_data['vertical_anchor']:
                    cell.vertical_anchor = eval(cell_data['vertical_anchor'])
                cell.text = ""
                paragraph = cell.text_frame.paragraphs[0]
                run = paragraph.add_run()
                run.text = cell_data['text']
                if cell_data['font_size']:
                    adjusted_size = cell_data['font_size'] * 0.8
                    run.font.size = Pt(adjusted_size)
                run.font.name = 'Arial'
                if cell_data.get('font_color'):
                    run.font.color.rgb = RGBColor.from_string(cell_data['font_color'])
                if 'bold' in cell_data:
                    run.font.bold = cell_data['bold']
                if 'italic' in cell_data:
                    run.font.italic = cell_data['italic']
                if cell_data['alignment']:
                    paragraph.alignment = get_alignment_value(cell_data['alignment'])
            except Exception as e:
                print(f"Error setting cell properties: {str(e)}")

def extract_text_from_slide(slide, slide_number, translate=False):
    """Extract all text elements from a slide."""
    slide_element = ET.Element("slide")
    slide_element.set("number", str(slide_number))
    for shape_index, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
            table_element = ET.SubElement(slide_element, "table_element")
            table_element.set("shape_index", str(shape_index))
            table_data = get_table_properties(shape.table)
            if translate:
                for row in table_data['cells']:
                    for cell in row:
                        cell['text'] = translate_text(cell['text'])
            props_element = ET.SubElement(table_element, "properties")
            props_element.text = json.dumps(table_data, indent=2)
        elif hasattr(shape, "text"):
            text_element = ET.SubElement(slide_element, "text_element")
            text_element.set("shape_index", str(shape_index))
            shape_data = get_shape_properties(shape)
            if translate:
                shape_data['text'] = translate_text(shape_data['text'])
            props_element = ET.SubElement(text_element, "properties")
            props_element.text = json.dumps(shape_data, indent=2)
    return slide_element

def ppt_to_xml(ppt_path: str, translate: bool = False) -> Optional[str]:
    """Convert PowerPoint to XML with intermediate saves."""
    root = ET.Element("presentation")
    base_dir = Path(ppt_path).parent
    try:
        prs = Presentation(ppt_path)
        root.set("file_path", os.path.basename(ppt_path))
        with ThreadPoolExecutor(max_workers=4) as executor:
            future_to_slide = {
                executor.submit(extract_text_from_slide, slide, slide_number, translate): slide_number 
                for slide_number, slide in enumerate(prs.slides, 1)
            }
            for future in future_to_slide:
                slide_number = future_to_slide[future]
                try:
                    slide_element = future.result()
                    root.append(slide_element)
                    intermediate_path = base_dir / f"slide_{slide_number}_{'translated' if translate else 'original'}.xml"
                    xml_str = minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
                    with open(intermediate_path, 'w', encoding='utf-8') as f:
                        f.write(xml_str)
                except Exception as e:
                    print(f"Error processing slide {slide_number}: {str(e)}")
        return minidom.parseString(ET.tostring(root)).toprettyxml(indent="  ")
    except Exception as e:
        print(f"Error processing presentation: {str(e)}")
        return None

def create_translated_ppt(original_ppt_path, translated_xml_path, output_ppt_path):
    """Create a new PowerPoint with translated text."""
    try:
        prs = Presentation(original_ppt_path)
        tree = ET.parse(translated_xml_path)
        root = tree.getroot()
        for slide_number, slide in enumerate(prs.slides, 1):
            xml_slide = root.find(f".//slide[@number='{slide_number}']")
            if xml_slide is None:
                continue
            for shape_index, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table_element = xml_slide.find(f".//table_element[@shape_index='{shape_index}']")
                    if table_element is not None:
                        props_element = table_element.find("properties")
                        if props_element is not None and props_element.text:
                            try:
                                table_data = json.loads(props_element.text)
                                apply_table_properties(shape.table, table_data)
                            except Exception as e:
                                print(f"Error applying table properties: {str(e)}")
                elif hasattr(shape, "text"):
                    text_element = xml_slide.find(f".//text_element[@shape_index='{shape_index}']")
                    if text_element is not None:
                        props_element = text_element.find("properties")
                        if props_element is not None and props_element.text:
                            try:
                                shape_data = json.loads(props_element.text)
                                apply_shape_properties(shape, shape_data)
                            except Exception as e:
                                print(f"Error applying shape properties: {str(e)}")
        prs.save(output_ppt_path)
        print(f"Translated PowerPoint saved to: {output_ppt_path}")
    except Exception as e:
        print(f"Error creating translated PowerPoint: {str(e)}")

def clean_path(path: str) -> str:
    """Remove quotes and handle escaped spaces in path."""
    path = path.strip("'\"")
    path = path.replace("\\ ", " ")
    path = path.replace("\\'", "'")
    return path

def translate_text_with_cache(text: str) -> str:
    """Translate text using a cache to avoid duplicate API calls."""
    if not text or text.isspace():
        return text
    with cache_lock:
        if text in translation_cache:
            return translation_cache[text]
    try:
        translated = translate_text(text)
        with cache_lock:
            translation_cache[text] = translated
        return translated
    except Exception as e:
        print(f"Translation error: {str(e)}")
        return text

def chunk_text(text: str, max_chunk_size: int = 1000) -> list[str]:
    """Split text into smaller chunks while preserving sentence boundaries."""
    if len(text) <= max_chunk_size:
        return [text]
    chunks = []
    current_chunk = []
    current_size = 0
    # naive sentence splitting
    sentences = text.replace('。', '.').replace('！', '!').replace('？', '?').split('.')
    for sentence in sentences:
        sentence = sentence.strip() + '.'
        if current_size + len(sentence) > max_chunk_size and current_chunk:
            chunks.append(''.join(current_chunk))
            current_chunk = []
            current_size = 0
        current_chunk.append(sentence)
        current_size += len(sentence)
    if current_chunk:
        chunks.append(''.join(current_chunk))
    return chunks

def cleanup_intermediate_files(base_dir: Path, pattern: str = "slide_*.xml"):
    """Clean up intermediate XML files after successful processing."""
    try:
        for file in base_dir.glob(pattern):
            file.unlink()
    except Exception as e:
        print(f"Warning: Could not clean up intermediate files: {str(e)}")

def process_ppt_file(ppt_path: Path, source_lang: str, target_lang: str):
    """Process a single PPT/PPTX file from XML extraction to final translation."""
    try:
        if not ppt_path.is_file():
            print(f"Error: '{ppt_path}' is not a valid file.")
            return
        if ppt_path.suffix.lower() not in ['.ppt', '.pptx']:
            print(f"Error: '{ppt_path}' is not a PowerPoint file.")
            return

        base_dir = ppt_path.parent

        # Adjust global translate_text to use chosen source/target
        global translate_text
        original_translate_func = translate_text
        translate_text = lambda text: original_translate_func(text, source_lang, target_lang)

        # Original XML
        print(f"Generating original XML for {ppt_path.name}...")
        original_xml = ppt_to_xml(str(ppt_path), translate=False)
        if original_xml:
            original_output_path = base_dir / f"{ppt_path.stem}_original.xml"
            with open(original_output_path, 'w', encoding='utf-8') as f:
                f.write(original_xml)
            print(f"Original XML saved: {original_output_path}")

        # Translated XML
        print(f"Generating translated XML (from {source_lang} to {target_lang}) for {ppt_path.name}...")
        translated_xml = ppt_to_xml(str(ppt_path), translate=True)
        if translated_xml:
            translated_output_path = base_dir / f"{ppt_path.stem}_translated.xml"
            with open(translated_output_path, 'w', encoding='utf-8') as f:
                f.write(translated_xml)
            print(f"Translated XML saved: {translated_output_path}")

            # Build final PPT
            print(f"Creating translated PPT for {ppt_path.name}...")
            output_filename = f"{ppt_path.stem}_translated{ppt_path.suffix}"
            output_ppt_path = base_dir / output_filename
            create_translated_ppt(str(ppt_path), str(translated_output_path), str(output_ppt_path))

            # Cleanup
            cleanup_intermediate_files(base_dir)
            print("Cleanup complete.")

        # Restore original translate_text
        translate_text = original_translate_func

    except Exception as e:
        print(f"Error in process_ppt_file for {ppt_path}: {str(e)}")

def main():
    try:
        path_input = input("Enter path to a PPTX file OR directory: ").strip()
        path_input = clean_path(path_input)
        path_input = os.path.expanduser(path_input)
        source_lang = input("Enter source language code (default 'zh'): ").strip().lower() or 'zh'
        target_lang = input("Enter target language code (default 'en'): ").strip().lower() or 'en'
        target_path = Path(path_input).resolve()

        if target_path.is_dir():
            # Recursively process all .ppt or .pptx
            for root, dirs, files in os.walk(target_path):
                for file in files:
                    if file.lower().endswith(('.ppt', '.pptx')):
                        full_path = Path(root) / file
                        process_ppt_file(full_path, source_lang, target_lang)
        else:
            # Process single file
            process_ppt_file(target_path, source_lang, target_lang)

    except Exception as e:
        print(f"Fatal error: {str(e)}")
        sys.exit(1)

if __name__ == "__main__":
    main()
